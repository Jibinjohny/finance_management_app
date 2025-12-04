#!/usr/bin/env python3
"""
Create PowerPoint presentation for CashFlow App
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("Installing python-pptx...")
    import subprocess
    subprocess.check_call(['pip3', 'install', 'python-pptx'])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    return slide

def add_content_slide(prs, title, content_points):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = title
    
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    for point in content_points:
        p = text_frame.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
    
    return slide

def add_section_header(prs, title):
    """Add a section header slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    title_shape = slide.shapes.title
    title_shape.text = title
    return slide

def create_presentation():
    """Create the main presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    add_title_slide(prs, 
                   "CashFlow App",
                   "Personal Finance Management System\\nVersion 1.0.0")
    
    # Overview
    add_content_slide(prs, "Project Overview", [
        "Modern personal finance management application",
        "Built with Flutter for cross-platform compatibility",
        "Local-first approach with SQLite database",
        "Comprehensive financial tracking and analytics",
        "User-friendly interface with glassmorphism design"
    ])
    
    # Key Features Section
    add_section_header(prs, "Key Features")
    
    add_content_slide(prs, "Account Management", [
        "Support for 9 account types",
        "Cash, Savings, Salary, Current, Credit Card",
        "Investment, Loan, Bank, and Other accounts",
        "Real-time balance tracking",
        "Color-coded visual identification"
    ])
    
    add_content_slide(prs, "Transaction Tracking", [
        "Income and expense recording",
        "Category-based organization",
        "Tag system for flexible filtering",
        "Multi-step transaction entry",
        "Search and filter capabilities"
    ])
    
    add_content_slide(prs, "Budgeting & Goals", [
        "Category-wise budget allocation",
        "Monthly and yearly budget periods",
        "Real-time progress tracking",
        "Overspending alerts",
        "Financial goal setting and monitoring"
    ])
    
    add_content_slide(prs, "Advanced Features", [
        "Recurring transaction automation",
        "EMI calculation and tracking",
        "Loan management with interest calculation",
        "Local push notifications",
        "Android home screen widget"
    ])
    
    add_content_slide(prs, "Analytics & Reports", [
        "Interactive dashboard with charts",
        "Net worth visualization",
        "Category-wise expense breakdown",
        "Monthly PDF report generation",
        "Account-wise financial trends"
    ])
    
    # Technical Architecture Section
    add_section_header(prs, "Technical Architecture")
    
    add_content_slide(prs, "Technology Stack", [
        "Framework: Flutter 3.9.2+",
        "Language: Dart",
        "Database: SQLite (Sqflite)",
        "State Management: Provider Pattern",
        "Charts: FL Chart library"
    ])
    
    add_content_slide(prs, "Architecture Pattern", [
        "Model-View-Provider (MVP)",
        "Clean Architecture principles",
        "Separation of concerns",
        "Reactive state management",
        "Repository pattern for data access"
    ])
    
    add_content_slide(prs, "Database Design", [
        "7 relational tables",
        "Users, Accounts, Transactions",
        "Budgets, Goals, Recurring Transactions",
        "Notifications and Tags",
        "Foreign key relationships for data integrity"
    ])
    
    add_content_slide(prs, "Security Features", [
        "Local-only data storage",
        "No cloud synchronization",
        "User authentication system",
        "Parameterized SQL queries",
        "OS-level app sandboxing"
    ])
    
    # User Experience Section
    add_section_header(prs, "User Experience")
    
    add_content_slide(prs, "Design Principles", [
        "Modern glassmorphism UI",
        "Dark mode support",
        "Smooth animations and transitions",
        "Intuitive navigation",
        "Responsive layouts"
    ])
    
    add_content_slide(prs, "Key Screens", [
        "Dashboard - Financial overview",
        "Accounts - Multi-account management",
        "Transactions - Income/expense tracking",
        "Statistics - Visual analytics",
        "Reports - Monthly summaries"
    ])
    
    # Implementation Highlights
    add_section_header(prs, "Implementation Highlights")
    
    add_content_slide(prs, "Smart Calculations", [
        "EMI calculation with interest rates",
        "Net worth computation across accounts",
        "Budget usage percentage tracking",
        "Loan balance and interest tracking",
        "Category-wise expense distribution"
    ])
    
    add_content_slide(prs, "Automation Features", [
        "Recurring transaction auto-creation",
        "EMI reminder scheduling (2 days before)",
        "Budget overspending notifications",
        "Goal deadline reminders",
        "Home widget auto-updates"
    ])
    
    add_content_slide(prs, "Data Management", [
        "SQLite database with version migrations",
        "Automatic balance updates",
        "Transaction history tracking",
        "PDF report generation and sharing",
        "Complete data reset option"
    ])
    
    # Platform Support
    add_section_header(prs, "Platform Support")
    
    add_content_slide(prs, "Supported Platforms", [
        "Android (5.0 and above)",
        "iOS (12.0 and above)",
        "Home screen widget (Android)",
        "Local notifications (Both platforms)",
        "Offline-first architecture"
    ])
    
    # Future Enhancements
    add_content_slide(prs, "Future Enhancements", [
        "Cloud backup and sync",
        "Multi-currency support",
        "Bank account integration",
        "Receipt scanning (OCR)",
        "Investment portfolio tracking",
        "Biometric authentication"
    ])
    
    # Statistics
    add_content_slide(prs, "Project Statistics", [
        "29 screens",
        "9 data models",
        "9 state providers",
        "4 core services",
        "14 reusable widgets",
        "Comprehensive test coverage"
    ])
    
    # Conclusion
    add_content_slide(prs, "Key Takeaways", [
        "Complete personal finance solution",
        "Privacy-focused with local storage",
        "Modern, intuitive user interface",
        "Robust technical architecture",
        "Scalable and maintainable codebase",
        "Ready for production deployment"
    ])
    
    # Thank You Slide
    add_title_slide(prs,
                   "Thank You",
                   "CashFlow App - Personal Finance Manager\\nVersion 1.0.0 | November 2025")
    
    return prs

def main():
    print("Creating PowerPoint presentation...")
    prs = create_presentation()
    
    output_file = "CashFlow_App_Presentation.pptx"
    prs.save(output_file)
    
    print(f"✅ Successfully created {output_file}")
    print(f"   Total slides: {len(prs.slides)}")
    
    import os
    size = os.path.getsize(output_file)
    print(f"   File size: {size / 1024:.1f} KB")

if __name__ == "__main__":
    main()
